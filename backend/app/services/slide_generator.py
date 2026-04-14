"""
Slide Generator - Servicio para generar y manipular diapositivas
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
import os


def extract_text_from_slides(pptx_path: str, up_to_slide: int = None, from_slide: int = 1) -> str:
    """
    Extrae texto de las diapositivas de un archivo PPTX

    Args:
        pptx_path: Ruta al archivo PPTX
        up_to_slide: Número de diapositiva hasta la cual extraer (1-indexed), None para todas
        from_slide: Número de diapositiva desde la cual comenzar (1-indexed), por defecto 1

    Returns:
        String con el texto extraído de las diapositivas
    """
    try:
        prs = Presentation(pptx_path)
        text_content = []

        # Determinar cuántas slides procesar
        total_slides = len(prs.slides)
        max_slide = up_to_slide if up_to_slide else total_slides
        start_slide = max(1, from_slide)  # Asegurar que sea al menos 1

        # Iterar sobre las slides usando índices (convertir de 1-indexed a 0-indexed)
        for idx in range(start_slide - 1, min(max_slide, total_slides)):
            slide = prs.slides[idx]
            slide_text = [f"--- Diapositiva {idx + 1} ---"]

            # Extraer texto de todas las formas en la diapositiva
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text.strip())

            text_content.append("\n".join(slide_text))

        return "\n\n".join(text_content)

    except Exception as e:
        print(f"Error extrayendo texto de PPTX: {e}")
        return "Contenido no disponible"


def create_example_slide(prs: Presentation, slide_data: dict) -> int:
    """
    Crea una diapositiva con estilo de ejemplo generado por IA

    Args:
        prs: Objeto Presentation de python-pptx
        slide_data: Dict con {title, content, key_points, is_analogy (opcional)}

    Returns:
        Índice de la nueva diapositiva
    """
    # Agregar diapositiva en blanco
    blank_layout = prs.slide_layouts[6]  # Layout en blanco
    slide = prs.slides.add_slide(blank_layout)

    # Dimensiones de la diapositiva
    slide_width = prs.slide_width
    slide_height = prs.slide_height

    # Detectar si es una analogía
    is_analogy = slide_data.get('is_analogy', False)

    # Agregar barra superior con indicador de IA
    ai_banner = slide.shapes.add_shape(
        1,  # Rectangle
        left=0,
        top=0,
        width=slide_width,
        height=Inches(0.6)
    )
    ai_banner.fill.solid()
    # Naranja para analogías, Púrpura para ejemplos normales
    banner_color = RGBColor(255, 152, 0) if is_analogy else RGBColor(103, 58, 183)
    ai_banner.fill.fore_color.rgb = banner_color
    ai_banner.line.color.rgb = banner_color

    # Texto del banner
    ai_text_frame = ai_banner.text_frame
    ai_text_frame.text = "Ejemplo con Analogía (IA)" if is_analogy else "Ejemplo Generado con IA"
    ai_text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    ai_text_frame.paragraphs[0].font.size = Pt(16)
    ai_text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    ai_text_frame.paragraphs[0].font.bold = True
    ai_text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Título del ejemplo
    title_box = slide.shapes.add_textbox(
        left=Inches(0.5),
        top=Inches(0.8),
        width=slide_width - Inches(1),
        height=Inches(0.8)
    )
    title_frame = title_box.text_frame
    title_frame.text = slide_data.get('title', 'Ejemplo')
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(33, 33, 33)

    # Contenido del ejemplo
    content_box = slide.shapes.add_textbox(
        left=Inches(0.5),
        top=Inches(1.8),
        width=slide_width - Inches(1),
        height=Inches(3.5)
    )
    content_frame = content_box.text_frame
    content_frame.word_wrap = True
    content_frame.text = slide_data.get('content', '')
    content_frame.paragraphs[0].font.size = Pt(18)
    content_frame.paragraphs[0].font.color.rgb = RGBColor(66, 66, 66)
    content_frame.paragraphs[0].space_after = Pt(12)

    # Puntos clave (si existen)
    key_points = slide_data.get('key_points', [])
    if key_points:
        points_y = Inches(5.5)

        # Título de puntos clave
        points_title_box = slide.shapes.add_textbox(
            left=Inches(0.5),
            top=points_y,
            width=slide_width - Inches(1),
            height=Inches(0.4)
        )
        points_title_frame = points_title_box.text_frame
        points_title_frame.text = "Puntos Clave:"
        points_title_frame.paragraphs[0].font.size = Pt(20)
        points_title_frame.paragraphs[0].font.bold = True
        points_title_frame.paragraphs[0].font.color.rgb = RGBColor(103, 58, 183)

        # Lista de puntos
        points_box = slide.shapes.add_textbox(
            left=Inches(0.7),
            top=points_y + Inches(0.5),
            width=slide_width - Inches(1.4),
            height=Inches(1.5)
        )
        points_frame = points_box.text_frame
        points_frame.word_wrap = True

        for i, point in enumerate(key_points[:4]):  # Máximo 4 puntos
            p = points_frame.add_paragraph() if i > 0 else points_frame.paragraphs[0]
            p.text = f"• {point}"
            p.font.size = Pt(14)
            p.font.color.rgb = RGBColor(66, 66, 66)
            p.space_before = Pt(6)

    # Agregar pequeño icono/texto de IA en la esquina inferior derecha
    footer_box = slide.shapes.add_textbox(
        left=slide_width - Inches(2.5),
        top=slide_height - Inches(0.4),
        width=Inches(2.3),
        height=Inches(0.3)
    )
    footer_frame = footer_box.text_frame
    footer_frame.text = "Generado con Gemini AI"
    footer_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT
    footer_frame.paragraphs[0].font.size = Pt(10)
    footer_frame.paragraphs[0].font.italic = True
    footer_frame.paragraphs[0].font.color.rgb = RGBColor(150, 150, 150)

    return len(prs.slides)


def insert_slide_in_presentation(pptx_path: str, slide_data: dict, position: int) -> int:
    """
    Inserta una diapositiva de ejemplo en una presentación existente

    Args:
        pptx_path: Ruta al archivo PPTX
        slide_data: Dict con datos del ejemplo
        position: Posición después de la cual insertar (0-indexed)

    Returns:
        Índice de la nueva diapositiva (1-indexed)
    """
    # Abrir presentación
    prs = Presentation(pptx_path)

    # Crear la diapositiva al final
    total_slides = create_example_slide(prs, slide_data)

    # Mover la diapositiva a la posición deseada
    # Las diapositivas en python-pptx están en prs.slides._sldIdLst
    # Para reordenar, necesitamos mover el XML
    slides = list(prs.slides._sldIdLst)

    # Índice de la nueva diapositiva (última)
    new_slide_idx = len(slides) - 1

    # Posición objetivo (después de 'position')
    target_idx = position  # position es 1-indexed desde el frontend, pero aquí es 0-indexed

    # Solo mover si no está ya en la posición correcta
    if new_slide_idx != target_idx:
        # Mover el elemento XML
        slide_element = slides[new_slide_idx]
        slides.insert(target_idx + 1, slides.pop(new_slide_idx))

    # Guardar presentación
    prs.save(pptx_path)

    return target_idx + 2  # Retornar posición 1-indexed (después de slide actual)


def create_question_slide(prs: Presentation, slide_data: dict) -> int:
    """
    Crea una diapositiva con estilo de pregunta generada por IA

    Args:
        prs: Objeto Presentation de python-pptx
        slide_data: Dict con {title, question, type, alternatives, correct_answer}

    Returns:
        Índice de la nueva diapositiva
    """
    # Agregar diapositiva en blanco
    blank_layout = prs.slide_layouts[6]  # Layout en blanco
    slide = prs.slides.add_slide(blank_layout)

    # Dimensiones de la diapositiva
    slide_width = prs.slide_width
    slide_height = prs.slide_height

    question_type = slide_data.get('type', 'open')

    # Agregar barra superior con indicador de IA (verde para preguntas)
    ai_banner = slide.shapes.add_shape(
        1,  # Rectangle
        left=0,
        top=0,
        width=slide_width,
        height=Inches(0.6)
    )
    ai_banner.fill.solid()
    ai_banner.fill.fore_color.rgb = RGBColor(46, 125, 50)  # Verde para preguntas
    ai_banner.line.color.rgb = RGBColor(46, 125, 50)

    # Texto del banner
    ai_text_frame = ai_banner.text_frame
    ai_text_frame.text = "Pregunta Generada con IA"
    ai_text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    ai_text_frame.paragraphs[0].font.size = Pt(16)
    ai_text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    ai_text_frame.paragraphs[0].font.bold = True
    ai_text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Título de la pregunta
    title_box = slide.shapes.add_textbox(
        left=Inches(0.5),
        top=Inches(0.8),
        width=slide_width - Inches(1),
        height=Inches(0.7)
    )
    title_frame = title_box.text_frame
    title_frame.text = slide_data.get('title', 'Pregunta de Evaluación')
    title_frame.paragraphs[0].font.size = Pt(28)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(33, 33, 33)

    # Pregunta principal
    question_box = slide.shapes.add_textbox(
        left=Inches(0.5),
        top=Inches(1.6),
        width=slide_width - Inches(1),
        height=Inches(1.5)
    )
    question_frame = question_box.text_frame
    question_frame.word_wrap = True
    question_frame.text = slide_data.get('question', '')
    question_frame.paragraphs[0].font.size = Pt(22)
    question_frame.paragraphs[0].font.bold = True
    question_frame.paragraphs[0].font.color.rgb = RGBColor(33, 150, 243)  # Azul para pregunta
    question_frame.paragraphs[0].space_after = Pt(16)

    # Si es de alternativas, mostrar opciones
    if question_type == "multiple-choice":
        alternatives = slide_data.get('alternatives', [])

        if alternatives:
            alternatives_y = Inches(3.3)

            # Contenedor de alternativas
            alternatives_box = slide.shapes.add_textbox(
                left=Inches(0.7),
                top=alternatives_y,
                width=slide_width - Inches(1.4),
                height=Inches(3.5)
            )
            alternatives_frame = alternatives_box.text_frame
            alternatives_frame.word_wrap = True

            for i, alt in enumerate(alternatives[:4]):  # Máximo 4 alternativas
                p = alternatives_frame.add_paragraph() if i > 0 else alternatives_frame.paragraphs[0]
                p.text = alt
                p.font.size = Pt(18)
                p.font.color.rgb = RGBColor(66, 66, 66)
                p.space_before = Pt(8)
                p.space_after = Pt(8)

            # Mostrar respuesta correcta (en texto pequeño al final)
            correct_answer = slide_data.get('correct_answer', '')
            if correct_answer:
                answer_box = slide.shapes.add_textbox(
                    left=Inches(0.5),
                    top=slide_height - Inches(0.7),
                    width=slide_width - Inches(1),
                    height=Inches(0.3)
                )
                answer_frame = answer_box.text_frame
                answer_frame.text = f"Respuesta correcta: {correct_answer}"
                answer_frame.paragraphs[0].font.size = Pt(14)
                answer_frame.paragraphs[0].font.color.rgb = RGBColor(46, 125, 50)
                answer_frame.paragraphs[0].font.bold = True

    else:  # Pregunta abierta
        # Agregar espacio para respuesta
        answer_space_box = slide.shapes.add_textbox(
            left=Inches(0.7),
            top=Inches(3.3),
            width=slide_width - Inches(1.4),
            height=Inches(3.0)
        )
        answer_space_frame = answer_space_box.text_frame
        answer_space_frame.text = "Espacio para desarrollo de respuesta"
        answer_space_frame.paragraphs[0].font.size = Pt(14)
        answer_space_frame.paragraphs[0].font.italic = True
        answer_space_frame.paragraphs[0].font.color.rgb = RGBColor(150, 150, 150)
        answer_space_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        answer_space_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Footer con indicador de IA
    footer_box = slide.shapes.add_textbox(
        left=slide_width - Inches(2.5),
        top=slide_height - Inches(0.4),
        width=Inches(2.3),
        height=Inches(0.3)
    )
    footer_frame = footer_box.text_frame
    footer_frame.text = "Generado con Gemini AI"
    footer_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT
    footer_frame.paragraphs[0].font.size = Pt(10)
    footer_frame.paragraphs[0].font.italic = True
    footer_frame.paragraphs[0].font.color.rgb = RGBColor(150, 150, 150)

    return len(prs.slides)


def insert_question_slide_in_presentation(pptx_path: str, slide_data: dict, position: int) -> int:
    """
    Inserta una diapositiva de pregunta en una presentación existente

    Args:
        pptx_path: Ruta al archivo PPTX
        slide_data: Dict con datos de la pregunta
        position: Posición después de la cual insertar (0-indexed)

    Returns:
        Índice de la nueva diapositiva (1-indexed)
    """
    # Abrir presentación
    prs = Presentation(pptx_path)

    # Crear la diapositiva al final
    total_slides = create_question_slide(prs, slide_data)

    # Mover la diapositiva a la posición deseada
    slides = list(prs.slides._sldIdLst)

    # Índice de la nueva diapositiva (última)
    new_slide_idx = len(slides) - 1

    # Posición objetivo (después de 'position')
    target_idx = position

    # Solo mover si no está ya en la posición correcta
    if new_slide_idx != target_idx:
        # Mover el elemento XML
        slide_element = slides[new_slide_idx]
        slides.insert(target_idx + 1, slides.pop(new_slide_idx))

    # Guardar presentación
    prs.save(pptx_path)

    return target_idx + 2  # Retornar posición 1-indexed (después de slide actual)


def create_analogy_slide(analogy_text: str, subject: str) -> str:
    """
    Crea una diapositiva con una analogía simple para explicar un concepto

    Args:
        analogy_text: Texto de la analogía
        subject: Asignatura para el título

    Returns:
        Ruta al archivo temporal de la diapositiva creada
    """
    import tempfile

    # Crear presentación temporal
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Agregar diapositiva en blanco
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    slide_width = prs.slide_width
    slide_height = prs.slide_height

    # Banner superior con color de analogía (naranja/dorado)
    ai_banner = slide.shapes.add_shape(
        1,  # Rectangle
        left=0,
        top=0,
        width=slide_width,
        height=Inches(0.6)
    )
    ai_banner.fill.solid()
    ai_banner.fill.fore_color.rgb = RGBColor(255, 152, 0)  # Naranja para analogías
    ai_banner.line.color.rgb = RGBColor(255, 152, 0)

    # Texto del banner
    ai_text_frame = ai_banner.text_frame
    ai_text_frame.text = "Ejemplo con Analogía (IA)"
    ai_text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    ai_text_frame.paragraphs[0].font.size = Pt(16)
    ai_text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    ai_text_frame.paragraphs[0].font.bold = True
    ai_text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Título
    title_box = slide.shapes.add_textbox(
        left=Inches(0.5),
        top=Inches(0.9),
        width=slide_width - Inches(1),
        height=Inches(0.7)
    )
    title_frame = title_box.text_frame
    title_frame.text = f"Explicación Simple - {subject}"
    title_frame.paragraphs[0].font.size = Pt(28)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(255, 152, 0)
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Contenedor para la analogía con fondo suave
    content_box = slide.shapes.add_shape(
        1,  # Rectangle
        left=Inches(0.5),
        top=Inches(1.8),
        width=slide_width - Inches(1),
        height=slide_height - Inches(2.3)
    )
    content_box.fill.solid()
    content_box.fill.fore_color.rgb = RGBColor(255, 248, 225)  # Amarillo muy suave
    content_box.line.color.rgb = RGBColor(255, 152, 0)
    content_box.line.width = Pt(2)

    # Texto de la analogía
    content_frame = content_box.text_frame
    content_frame.word_wrap = True
    content_frame.margin_left = Inches(0.3)
    content_frame.margin_right = Inches(0.3)
    content_frame.margin_top = Inches(0.3)
    content_frame.margin_bottom = Inches(0.3)

    p = content_frame.paragraphs[0]
    p.text = analogy_text
    p.font.size = Pt(20)
    p.font.color.rgb = RGBColor(51, 51, 51)
    p.line_spacing = 1.3
    p.space_before = Pt(12)
    p.space_after = Pt(12)

    # Guardar en archivo temporal
    temp_file = tempfile.NamedTemporaryFile(delete=False, suffix='.pptx')
    prs.save(temp_file.name)
    temp_file.close()

    return temp_file.name
